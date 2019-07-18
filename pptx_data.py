import requests
import browser_cookie3

cj = browser_cookie3.chrome()

fname = requests.get('URL', cookies = cj)


file = open('file.pptx', 'wb')
file.write(fname.content)
file.close()

fname1 = 'file.pptx'

import pptx
from pptx import Presentation


ppt = Presentation(fname1)

text_runs = []

# for slide in ppt.slides:
#     for shape in slide.shapes:
#         text_runs.append(shape.text)

for slide in ppt.slides:
    for shape in slide.shapes:
        if(shape.has_text_frame):
            text_runs.append(shape.text)


import re

urls = re.compile('http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+')
urllist = urls.findall(str(text_runs))


for item in urllist:
    if (len(str(item))< 1):
        urllist.remove(str(item))


#print(urllist)
import zipfile

zipref = zipfile.ZipFile(fname1)
zipref.extractall('fname1')

import os
path = 'fname1/ppt/slides/_rels'

slidename = []

for filename in os.listdir(path):
    slidename.append(filename)

#print(slidename)
content_list = []
for i in range(slidename.__len__()):
    path = path+'/'+slidename[i]
    slidedata = open(path, 'r')
    path = 'fname1/ppt/slides/_rels'
    content_list.append(slidedata.read())
    slidedata.close()


#print(content_list)
xml_urls = []
for i in range(content_list.__len__()):
    xml_urls.append(urls.findall(content_list[i]))

#print(xml_urls)

removelist = set(['http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
 'http://schemas.openxmlformats.org/package/2006/relationships',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/printerSettings',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing',
 'http://schemas.openxmlformats.org/package/2006/relationships',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/settings',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme',
 'http://schemas.microsoft.com/office/2007/relationships/stylesWithEffects',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/fontTable',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/webSettings'
 'http://schemas.openxmlformats.org/package/2006/relationships',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout',
 'http://schemas.openxmlformats.org/package/2006/relationships',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/webSettings',
 'http://schemas.microsoft.com/sharepoint/v3/contenttype/forms',
 'http://schemas.microsoft.com/office/2006/metadata/properties/metaAttributes',
 'http://purl.org/dc/elements/1.1/',
 'http://schemas.microsoft.com/office/2006/metadata/longProperties',
 'http://purl.org/dc/terms/',
 'http://www.w3.org/2001/XMLSchema-instance',
 'http://schemas.microsoft.com/office/2006/metadata/properties',
 'http://dublincore.org/schemas/xmls/qdc/2003/04/02/dcterms.xsd',
 'http://schemas.microsoft.com/office/infopath/2007/PartnerControls',
 'http://schemas.microsoft.com/office/2006/documentManagement/types',
 'http://www.w3.org/2001/XMLSchema',
 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
 'http://schemas.openxmlformats.org/officeDocument/2006/customXml',
 'http://schemas.microsoft.com/internal/obd',
 'http://schemas.openxmlformats.org/drawingml/2006/main',
 'http://dublincore.org/schemas/xmls/qdc/2003/04/02/dc.xsd',
 'http://schemas.microsoft.com/office/2006/metadata/contentType',
 'http://schemas.microsoft.com/office/infopath/2007/Partner<?xml',
 'http://schemas.microsoft.com/sharepoint/events',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/header',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXml',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/package',
 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/endnotes',
 'http://www.w3.org/1999/02/22-rdf-syntax-ns',
 'http://schemas.openxmlformats.org/officeDocument/2006/bibliography'])

clean_xml = []

for i in xml_urls:
    clean_xml += i

total_url = set(clean_xml + urllist) -set(removelist)

print(total_url)